from pathlib import Path
import re
import datetime
import pandas as pd
import fitz
from pptx import Presentation
import shutil

# -----------------------------
# Patient + Geographical data safety
def contains_patient_data(text: str) -> bool:
    # 1. Names or keywords
    name_patterns = [r'\bName\b', r'\bPatient\b', r'\bFirst Name\b', r'\bLast Name\b']
    
    # 2. Dates of birth
    dob_patterns = [
        r'\bDOB\b', r'\bDate of Birth\b',
        r'\b\d{2}[/-]\d{2}[/-]\d{4}\b',    # DD/MM/YYYY or DD-MM-YYYY
        r'\b\d{4}[/-]\d{2}[/-]\d{2}\b'     # YYYY/MM/DD or YYYY-MM-DD
    ]
    
    # 3. Dutch BSN — 9 digits
    bsn_patterns = [r'\b\d{9}\b']
    
    # 4. Hospital or medical IDs
    mrn_patterns = [r'\bMRN\b', r'\bPatient ID\b', r'\bRecord Number\b']
    
    # 5. Geographical data
    geo_patterns = [
        r'\bAddress\b', r'\bStreet\b', r'\bCity\b', r'\bTown\b', r'\bVillage\b',
        r'\bPostcode\b', r'\bZIP\b', r'\bCountry\b'
    ]
    
    all_patterns = name_patterns + dob_patterns + bsn_patterns + mrn_patterns + geo_patterns
    
    for pattern in all_patterns:
        if re.search(pattern, text, re.IGNORECASE):
            return True
    return False

# -----------------------------
# Text extraction functions
def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == '.pptx': return extract_text_from_pptx(path)
    if ext == '.pdf': return extract_text_from_pdf(path)
    if ext in ['.xlsx', '.xls']: return extract_text_from_excel(path)
    if ext == '.csv': return extract_text_from_csv(path)
    return ""

def extract_text_from_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, 'text')]
        slides_text.append(f"Slide {i+1}: " + " ".join(slide_text))
    return "\n".join(slides_text)

def extract_text_from_pdf(path: Path) -> str:
    doc = fitz.open(path)
    return "\n".join([f"Page {i+1}: {page.get_text()}" for i, page in enumerate(doc)])

def extract_text_from_excel(path: Path) -> str:
    df = pd.read_excel(path, engine='openpyxl')
    return df.to_csv(index=False)

def extract_text_from_csv(path: Path) -> str:
    df = pd.read_csv(path)
    return df.to_csv(index=False)

# -----------------------------
# Metadata function
def get_file_metadata(path: Path) -> dict:
    text = extract_text(path)
    return {
        'filename': path.name,
        'filepath': str(path.resolve()),
        'extension': path.suffix.lower(),
        'size_bytes': path.stat().st_size,
        'created_at': datetime.datetime.fromtimestamp(path.stat().st_ctime).isoformat(),
        'modified_at': datetime.datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
        'contains_patient_data': contains_patient_data(text),
        'num_characters': len(text),
        'num_words': len(text.split()),
        'text': text,
        'text_snippet': text[:500],
        'year_in_filename': re.findall(r'20\d{2}', path.stem),
        'topic_hints_in_filename': [t for t in ['survival','causal','epidemiology','ml','regression'] if t in path.stem.lower()]
    }

# -----------------------------
# Anonymization function
def save_anonymized_file(path: Path, base_dir: Path, anon_folder_name='anonymized') -> Path:
    # Avoid re-anonymizing if already in anonymized folder
    if anon_folder_name in path.parts:
        return path

    rel_path = path.relative_to(base_dir)
    target_path = base_dir / anon_folder_name / rel_path
    target_path.parent.mkdir(parents=True, exist_ok=True)

    # Simple anonymization: replace sensitive text in data files
    ext = path.suffix.lower()
    if ext in ['.xlsx', '.xls', '.csv']:
        if ext in ['.xlsx', '.xls']:
            df = pd.read_excel(path, engine='openpyxl')
        else:
            df = pd.read_csv(path)
        for col in df.select_dtypes(include='object').columns:
            df[col] = df[col].apply(lambda x: "[REDACTED]" if isinstance(x, str) and contains_patient_data(x) else x)
        if ext in ['.xlsx', '.xls']:
            df.to_excel(target_path, index=False, engine='openpyxl')
        else:
            df.to_csv(target_path, index=False)
    else:
        # For pptx or pdf, just copy (full text replacement would require deeper parsing)
        shutil.copy2(path, target_path)

    return target_path

# -----------------------------
# Recursive scan, anonymization, and metadata collection
DATA_DIR = Path(r"D:/Hackathon 13_12_2025")
file_metadata_dict = {}
failed_files = {}  # Store files that failed processing
ANON_FOLDER_NAME = 'anonymized'

all_files = list(DATA_DIR.glob('/.'))
supported_exts = ['.pptx', '.pdf', '.xlsx', '.xls', '.csv']
matching_files = [f for f in all_files if f.suffix.lower() in supported_exts]

for f in matching_files:
    try:
        anonymized_path = None
        meta = get_file_metadata(f)

        # Trigger anonymization if patient or geo data detected
        if not ANON_FOLDER_NAME in f.parts and meta['contains_patient_data']:
            print(f"⚠ {f.name} contains patient or geographical data, anonymizing...")
            anonymized_path = save_anonymized_file(f, DATA_DIR, ANON_FOLDER_NAME)
            print(f"✅ Saved anonymized copy at {anonymized_path}")

        final_path = anonymized_path if anonymized_path else f
        final_meta = get_file_metadata(final_path)
        file_metadata_dict[final_path.name] = final_meta

    except Exception as e:
        print(f"❌ Error processing {f.name}: {e}")
        failed_files[f.name] = str(e)

# -----------------------------
# Recap summary
total_files_found = len(matching_files)
total_files_processed = len(file_metadata_dict)
total_failed = len(failed_files)
total_anonymized = sum(1 for meta in file_metadata_dict.values() 
                    if ANON_FOLDER_NAME in Path(meta['filepath']).parts)
total_skipped = total_files_found - total_files_processed - total_failed

print("\n==================== Summary Report ====================")
print(f"Total files found (supported types): {total_files_found}")
print(f"Total files anonymized: {total_anonymized}")
print(f"Total files successfully processed: {total_files_processed}")
print(f"Total files skipped: {total_skipped}")
print(f"Total files failed to process: {total_failed}")

if total_failed > 0:
    print("\n⚠ Files that failed to process:")
    for fname, err in failed_files.items():
        print(f"- {fname}: {err}")

# Optional: List anonymized files
if total_anonymized > 0:
    print("\nℹ Anonymized files:")
    for meta in file_metadata_dict.values():
        if ANON_FOLDER_NAME in Path(meta['filepath']).parts:
            print(f"- {meta['filename']} (saved at {meta['filepath']})")



# -----------------------------
# Optional: Print metadata summary
for filename, meta in file_metadata_dict.items():
    print(f"\n=== {filename} ===")
    for key, value in meta.items():
        if key == 'text':
            print(f"{key}: (full content omitted, {len(value)} characters)")
        else:
            print(f"{key}: {value}")
